# scripts/make_presentation.py
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from pathlib import Path

OUTPUT_DIR = Path('output')
OUT_PPTX = Path('docs/Acquisition_Insights.pptx')

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_funnel_slide(prs, funnel_img):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Acquisition Funnel (Visits → Starts)"
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6)).text_frame
    tf.text = "Funnel by channel. Visits are aggregated visit_count; Starts are unique application_id."
    if Path(funnel_img).exists():
        slide.shapes.add_picture(str(funnel_img), Inches(0.5), Inches(1.2), width=Inches(9))

def add_impact_slide(prs, lost_csv, sim_csv, chart_img=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Impact & Recommendations"
    tf = slide.placeholders[1].text_frame
    lost = pd.read_csv(lost_csv)
    sim = pd.read_csv(sim_csv)
    total_lost = lost['lost_revenue'].sum()
    total_inc_revenue = sim['incremental_revenue'].sum()
    tf.text = f"Sample assumptions: avg_rev £200; approval_rate 60%; +5pp Visits→Starts uplift modelled."
    p = tf.add_paragraph(); p.text = f"Estimated lost revenue (abandonments): £{total_lost:,.2f}"
    p = tf.add_paragraph(); p.text = f"Estimated incremental revenue from +5pp uplift: £{total_inc_revenue:,.2f}"
    p = tf.add_paragraph(); p.text = "Recommendations:"
    p = tf.add_paragraph(); p.text = "1) Provide Submissions and Decisions (with revenue) and marketing spend for ROI analysis."
    p = tf.add_paragraph(); p.text = "2) Prioritise App UX tests (highest absolute traffic)."
    if chart_img and Path(chart_img).exists():
        slide.shapes.add_picture(str(chart_img), Inches(6.5), Inches(1.2), width=Inches(3))

def make_presentation():
    prs = Presentation()
    add_title_slide(prs, "Credit Card Acquisition — Funnel & Commerce Impact", "Sample analysis; assumptions in notes")
    add_funnel_slide(prs, 'output/funnel_chart.png')
    add_impact_slide(prs, 'output/lost_estimates_by_channel.csv', 'output/5pp_uplift_sim_by_channel.csv', 'output/channel_starts.png')
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print("Saved PPTX to", OUT_PPTX)

if __name__ == '__main__':
    make_presentation()
